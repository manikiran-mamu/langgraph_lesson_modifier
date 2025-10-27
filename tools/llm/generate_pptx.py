# tools/llm/generate_pptx.py

import os
import uuid
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

FONT_NAME = "Poppins"

def _add_title_slide(prs, title_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.name = FONT_NAME

def _add_center_text_slide(prs, text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(32)
    p.font.name = FONT_NAME
    p.alignment = 1  # center

def _add_content_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.name = FONT_NAME

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for line in content.split("\n"):
        if line.strip():
            p = tf.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(20)
            p.font.name = FONT_NAME

def generate_slide_deck(sections: dict, lesson_objective: str, language_objective: str) -> str:
    prs = Presentation()

    # Slide 1: Lesson Title
    _add_title_slide(prs, "Lesson Plan Slide Deck")

    # Slide 2: Engager - Intro Teacher Activities
    _add_content_slide(prs, "Engager", sections.get("intro_teacher", "No content"))

    # Slide 3: Content (Lesson + Language Objective)
    combined_objectives = f"Lesson Objective:\n{lesson_objective}\n\nLanguage Objective:\n{language_objective}"
    _add_content_slide(prs, "Objectives", combined_objectives)

    # Slide 4: I DO Divider
    _add_center_text_slide(prs, "I DO")

    # Slides 5+: Based on I DO + Lesson Content
    i_do_text = sections.get("i_do_teacher", "")
    lesson_content = sections.get("content", "")
    combined = f"{lesson_objective}\n\n{i_do_text}\n\n{lesson_content}"
    for chunk in combined.split("\n\n"):
        _add_content_slide(prs, "Direct Instruction", chunk.strip())

    # Slide: WE DO Divider
    _add_center_text_slide(prs, "WE DO")

    # Slides: 2 slides from we_do_teacher and lesson_objective
    we_do_text = sections.get("we_do_teacher", "")
    content_chunks = [lesson_objective, we_do_text]
    for chunk in content_chunks:
        _add_content_slide(prs, "Guided Practice", chunk.strip())

    # Save file
    save_dir = "data/outputs/slides"
    os.makedirs(save_dir, exist_ok=True)
    filename = f"lesson_slides_{uuid.uuid4().hex}.pptx"
    path = os.path.join(save_dir, filename)
    prs.save(path)
    return path

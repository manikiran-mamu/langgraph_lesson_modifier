from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os
import uuid


# -----------------------------------------------------------
# Helper Function: Format content with advanced rules
# -----------------------------------------------------------
def _set_formatted_content(text_frame, content, center=False):
    """
    Handles 2 main cases:
    1️⃣ Only '\n\n' present → first part Blue, second part Red.
    2️⃣ Both '\n' and '\n\n' present →
        - Before first '\n'        → Blue
        - Between '\n' and '\n\n'  → Blue Bold
        - Between '\n\n' and next '\n' → Red
        - After that (final)       → Red Bold
    """
    text_frame.clear()
    text_frame.word_wrap = True
    if center:
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Case 1: Has both \n and \n\n
    if "\n\n" in content and "\n" in content:
        # Find positions of key delimiters
        first_single = content.find("\n")
        first_double = content.find("\n\n")
        second_single = content.find("\n", first_double + 2)

        # Segment extraction with safety
        part1 = content[:first_single].strip() if first_single != -1 else ""
        part2 = content[first_single:first_double].strip() if first_double != -1 else ""
        part3 = content[first_double:second_single].replace("\n\n", "").strip() if second_single != -1 else content[first_double:].replace("\n\n", "").strip()
        part4 = content[second_single:].strip() if second_single != -1 else ""

        styled_parts = [
            (part1, RGBColor(0, 102, 204), False),  # 🔵 Blue
            (part2, RGBColor(0, 102, 204), True), # 🔵 Blue Bold
            ("__SPACER__", None, None),       
            (part3, RGBColor(255, 0, 0), False),    # 🔴 Red
            (part4, RGBColor(255, 0, 0), True),     # 🔴 Red Bold
        ]

    # Case 2: Only '\n\n' present
    elif "\n\n" in content:
        english_part, translation_part = content.split("\n\n", 1)
        styled_parts = [
            (english_part.strip(), RGBColor(0, 102, 204), False),  # 🔵 Blue
            (translation_part.strip(), RGBColor(255, 0, 0), False)  # 🔴 Red
        ]

    # Case 3: Fallback → treat all as blue
    else:
        styled_parts = [(content.strip(), RGBColor(0, 0, 0), False)]

    # Render styled parts
    for text, color, bold in styled_parts:
        if not text:
            continue
        for line in text.split("\n"):
            clean_line = line.strip()
            if not clean_line:
                continue
            p = text_frame.add_paragraph()
            p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = clean_line
            run.font.size = Pt(16)
            run.font.name = "Poppins"
            run.font.color.rgb = color
            run.font.bold = bold


# -----------------------------------------------------------
# Layout Functions
# -----------------------------------------------------------
def _add_title_content_slide(prs, title, content):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Poppins"
    run.font.bold = True
    run.font.size = Pt(22)
    p.alignment = PP_ALIGN.LEFT

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    _set_formatted_content(text_frame, content)


def _add_title_only_slide(prs, title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = title_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Poppins"
    run.font.bold = True
    run.font.size = Pt(36)
    p.alignment = PP_ALIGN.CENTER


def _add_content_only_slide(prs, content):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_formatted_content(text_frame, content, center=True)


# -----------------------------------------------------------
# Main Deck Generator
# -----------------------------------------------------------
def generate_slide_deck(slides: list) -> str:
    prs = Presentation()

    for slide in slides:
        title = slide.get("title", "").strip()
        content = slide.get("content", "").strip()

        if not content:
            _add_title_only_slide(prs, title)
        elif not title:
            _add_content_only_slide(prs, content)
        elif title.lower() in ["i do", "we do", "you do", "engager"]:
            _add_title_only_slide(prs, title)
        else:
            _add_title_content_slide(prs, title, content)

    save_dir = "data/outputs/slides"
    os.makedirs(save_dir, exist_ok=True)
    filename = f"lesson_slides_{uuid.uuid4().hex}.pptx"
    path = os.path.join(save_dir, filename)
    prs.save(path)
    return path